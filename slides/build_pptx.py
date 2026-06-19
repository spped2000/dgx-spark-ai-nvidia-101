# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0"]
# ///
"""
build_pptx.py — สร้าง deck.pptx (เปิดใน PowerPoint/Google Slides/Keynote) จาก slides.json
รัน: uv run slides/build_pptx.py   (จากโฟลเดอร์ course/)  หรือ  uv run build_pptx.py (จาก slides/)
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "slides.json")
OUT = os.path.join(HERE, "deck.pptx")
FONT = "Tahoma"  # รองรับทั้งไทยและอังกฤษ
NV_GREEN = RGBColor(0x76, 0xB9, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)

TITLE = "DGX Spark — AI NVIDIA 101 / Agent Enterprise"
SUBTITLE = "AGICAFET · 20 มิ.ย. 2026 · 10:00–16:00 · Cleverse Rama9 · 30 ที่นั่ง · 100% no-code"


def set_font(tf, size=None, bold=None, color=None):
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = FONT
            if size: r.font.size = Pt(size)
            if bold is not None: r.font.bold = bold
            if color: r.font.color.rgb = color


def add_title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(2))
    tf = box.text_frame; tf.word_wrap = True
    tf.text = TITLE
    set_font(tf, size=40, bold=True, color=DARK)
    b2 = s.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(12), Inches(1))
    tf2 = b2.text_frame; tf2.word_wrap = True; tf2.text = SUBTITLE
    set_font(tf2, size=18, color=NV_GREEN)
    return s


def add_section_slide(prs, title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bar = s.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(12), Inches(1.5))
    tf = bar.text_frame; tf.word_wrap = True; tf.text = title
    set_font(tf, size=30, bold=True, color=NV_GREEN)
    return s


def add_content_slide(prs, title, bullets, notes=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # title
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True; tf.text = title
    set_font(tf, size=26, bold=True, color=DARK)
    # accent line via a thin textbox underline (simple)
    # bullets
    bb = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.4))
    tf2 = bb.text_frame; tf2.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = "•  " + str(b)
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(18); r.font.color.rgb = DARK
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def main():
    if not os.path.exists(SRC):
        print(f"[FATAL] ไม่พบ {SRC} — รัน slide workflow ให้เสร็จก่อน", file=sys.stderr); sys.exit(2)
    data = json.load(open(SRC, encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)  # 16:9

    add_title_slide(prs)
    # agenda
    ag = add_content_slide(prs, "Agenda — 12 Modules", [m.get("module", "") for m in data],
                           notes="ภาพรวมทั้งวัน 10:00–16:00")
    n = 0
    for m in data:
        add_section_slide(prs, m.get("module", ""))
        for sl in m.get("slides", []):
            add_content_slide(prs, sl.get("title", ""), sl.get("bullets", []), sl.get("notes", ""))
            n += 1
    # thanks
    add_section_slide(prs, "ขอบคุณครับ / Q&A — AGICAFET")
    total = len(prs.slides._sldIdLst)
    prs.save(OUT)
    print(f"[done] เขียน {OUT}  ({total} slides รวม, {n} content slides, {len(data)} modules)")


if __name__ == "__main__":
    main()
